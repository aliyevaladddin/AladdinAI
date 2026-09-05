# NOTICE: This file is protected under RCF-PL v2.0.3
"""PPTX format converter — bidirectional .pptx ↔ .wrt.

Converts PowerPoint presentations to/from .wrt format for AI editing.
Slides are represented as numbered sections with headings and content.

Format in .wrt:
    [slide 1]
    [h1]Title[/h1]
    Slide content...
    [/slide]

    [slide 2]
    [h2]Subtitle[/h2]
    More content...
    [/slide]
"""
from __future__ import annotations

import io
import logging
import re

log = logging.getLogger(__name__)


def pptx_to_wrt(pptx_bytes: bytes) -> str:
    """Convert a .pptx file (bytes) to .wrt text format."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    lines: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"[slide {slide_num}]")

        # Extract title if present
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                lines.append(f"[h1]{title_text}[/h1]")

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    lines.append(text)

            # Handle tables
            if shape.shape_type == 19:  # Table
                lines.append("[table]")
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                lines.append("[/table]")

        lines.append("[/slide]")
        lines.append("")  # Blank line between slides

    return "\n".join(lines).strip() + "\n"


def wrt_to_pptx(wrt_text: str) -> bytes:
    """Convert .wrt text format to a .pptx file (bytes)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Parse slides
    slide_pattern = re.compile(r'\[slide \d+\](.*?)\[/slide\]', re.DOTALL)
    slides_content = slide_pattern.findall(wrt_text)

    for slide_text in slides_content:
        # Use blank layout (6) for custom content
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Extract title
        title_match = re.search(r'\[h1\](.*?)\[/h1\]', slide_text)
        if title_match:
            title_text = title_match.group(1).strip()
            # Add title textbox
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title_text
            p = title_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            # Remove title from content
            slide_text = slide_text.replace(title_match.group(0), "")

        # Extract content (everything except title and table)
        content_lines = []
        in_table = False
        for line in slide_text.split("\n"):
            stripped = line.strip()
            if stripped == "[table]":
                in_table = True
                continue
            if stripped == "[/table]":
                in_table = False
                continue
            if not in_table and stripped and not stripped.startswith("["):
                content_lines.append(stripped)

        # Add content textbox if there's content
        if content_lines:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_box.text_frame
            content_frame.text = "\n".join(content_lines)
            p = content_frame.paragraphs[0]
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
